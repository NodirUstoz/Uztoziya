from rest_framework import generics, status, permissions
from rest_framework.decorators import api_view, permission_classes
from rest_framework.response import Response
from rest_framework.permissions import IsAuthenticated
from django.shortcuts import get_object_or_404, render
from django.http import HttpResponse, Http404, FileResponse, JsonResponse
from django.db.models import Q
from django.utils import timezone
from django.utils.text import slugify
from django.contrib.auth.decorators import login_required
from django.conf import settings
from django.urls import reverse
from django.views.decorators.csrf import ensure_csrf_cookie
from django.views.decorators.http import require_POST
import io
import logging
import os
import re
import json
from openai import OpenAI, OpenAIError
from pptx import Presentation

try:
    from google import genai
except ImportError:  # pragma: no cover - optional dependency
    genai = None


logger = logging.getLogger(__name__)


TEXT_PROVIDER_LABELS = {
    'openai': "OpenAI GPT-4o",
    'openrouter': "OpenRouter GPT-4o",
    'gemini': "Google Gemini 1.5",
}

IMAGE_PROVIDER_LABELS = {
    'openai': "OpenAI GPT-4o tasvir modeli",
    'gemini': "Google Gemini 1.5",
}


def _resolve_text_provider() -> str:
    if getattr(settings, 'OPENAI_API_KEY', ''):
        return 'openai'
    if getattr(settings, 'OPENROUTER_API_KEY', ''):
        return 'openrouter'
    if getattr(settings, 'GOOGLE_GEMINI_API_KEY', ''):
        return 'gemini'
    raise ValueError("Hech qanday matnli AI API kaliti topilmadi. Iltimos, .env faylini tekshiring.")


def _resolve_image_provider() -> str:
    if getattr(settings, 'OPENAI_API_KEY', ''):
        return 'openai'
    if getattr(settings, 'GOOGLE_GEMINI_API_KEY', ''):
        return 'gemini'
    raise ValueError("Tasvir yaratish uchun OpenAI yoki Google Gemini API kaliti talab qilinadi.")


def _get_provider_label(provider: str, mapping: dict[str, str]) -> str:
    return mapping.get(provider, provider.title())

from .models import (
    Material, MaterialCategory, MaterialRating, MaterialDownload,
    Assignment, StudentSubmission, VideoLesson, Model3D
)
from .serializers import (
    MaterialCategorySerializer,
    MaterialSerializer,
    MaterialRatingSerializer,
    AssignmentSerializer,
    StudentSubmissionSerializer,
    VideoLessonSerializer,
    Model3DSerializer
)


class MaterialCategoryListView(generics.ListAPIView):
    """Material kategoriyalari ro'yxati"""
    queryset = MaterialCategory.objects.all()
    serializer_class = MaterialCategorySerializer
    permission_classes = [permissions.AllowAny]


class MaterialListView(generics.ListCreateAPIView):
    """Materiallar ro'yxati va yaratish"""
    serializer_class = MaterialSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = Material.objects.filter(is_public=True)
        
        # Filtrlash
        category = self.request.query_params.get('category')
        material_type = self.request.query_params.get('type')
        subject = self.request.query_params.get('subject')
        grade_level = self.request.query_params.get('grade')
        search = self.request.query_params.get('search')
        
        if category:
            queryset = queryset.filter(category_id=category)
        if material_type:
            queryset = queryset.filter(material_type=material_type)
        if subject:
            queryset = queryset.filter(author__subject=subject)
        if grade_level:
            queryset = queryset.filter(grade_level=grade_level)
        if search:
            queryset = queryset.filter(
                Q(title__icontains=search) |
                Q(description__icontains=search) |
                Q(tags__icontains=search)
            )
        
        return queryset.order_by('-created_at')
    
    def perform_create(self, serializer):
        serializer.save(author=self.request.user)


class MaterialDetailView(generics.RetrieveUpdateDestroyAPIView):
    """Material tafsilotlari"""
    serializer_class = MaterialSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        return Material.objects.filter(
            Q(is_public=True) | Q(author=self.request.user)
        )


class MaterialCreateView(generics.CreateAPIView):
    """Yangi material yaratish"""
    serializer_class = MaterialSerializer
    permission_classes = [IsAuthenticated]
    
    def perform_create(self, serializer):
        serializer.save(author=self.request.user)


class MaterialUpdateView(generics.UpdateAPIView):
    """Materialni yangilash"""
    serializer_class = MaterialSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        return Material.objects.filter(author=self.request.user)


class MaterialDeleteView(generics.DestroyAPIView):
    """Materialni o'chirish"""
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        return Material.objects.filter(author=self.request.user)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def download_material(request, pk):
    """Materialni yuklab olish"""
    material = get_object_or_404(Material, pk=pk)
    
    # Yuklab olish huquqini tekshirish
    if not material.is_public and material.author != request.user:
        return Response({
            'error': 'Bu materialni yuklab olish huquqingiz yo\'q'
        }, status=status.HTTP_403_FORBIDDEN)
    
    try:
        # Yuklab olish statistikasini yangilash
        material.download_count += 1
        material.save()
        
        # Yuklab olish tarixini saqlash
        MaterialDownload.objects.create(
            material=material,
            user=request.user,
            ip_address=request.META.get('REMOTE_ADDR')
        )
        
        # Faylni yuklab olish
        if os.path.exists(material.file.path):
            with open(material.file.path, 'rb') as f:
                response = HttpResponse(f.read(), content_type='application/octet-stream')
                response['Content-Disposition'] = f'attachment; filename="{material.title}.{material.file.name.split(".")[-1]}"'
                return response
        else:
            return Response({
                'error': 'Fayl topilmadi'
            }, status=status.HTTP_404_NOT_FOUND)
            
    except Exception as e:
        return Response({
            'error': 'Faylni yuklab olishda xatolik'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def rate_material(request, pk):
    """Materialni baholash"""
    material = get_object_or_404(Material, pk=pk)
    
    rating_value = request.data.get('rating')
    comment = request.data.get('comment', '')
    
    if not rating_value or not (1 <= int(rating_value) <= 5):
        return Response({
            'error': 'Reyting 1 dan 5 gacha bo\'lishi kerak'
        }, status=status.HTTP_400_BAD_REQUEST)
    
    # Mavjud reytingni yangilash yoki yangi yaratish
    rating, created = MaterialRating.objects.get_or_create(
        material=material,
        user=request.user,
        defaults={'rating': rating_value, 'comment': comment}
    )
    
    if not created:
        rating.rating = rating_value
        rating.comment = comment
        rating.save()
    
    # Material reytingini yangilash
    ratings = MaterialRating.objects.filter(material=material)
    avg_rating = sum(r.rating for r in ratings) / ratings.count()
    material.rating = avg_rating
    material.save()
    
    return Response({
        'message': 'Reyting muvaffaqiyatli saqlandi',
        'rating': MaterialRatingSerializer(rating).data
    })


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def search_materials(request):
    """Materiallarni qidirish"""
    query = request.query_params.get('q', '')
    category = request.query_params.get('category')
    material_type = request.query_params.get('type')
    subject = request.query_params.get('subject')
    grade_level = request.query_params.get('grade')
    
    queryset = Material.objects.filter(is_public=True)
    
    if query:
        queryset = queryset.filter(
            Q(title__icontains=query) |
            Q(description__icontains=query) |
            Q(tags__icontains=query)
        )
    
    if category:
        queryset = queryset.filter(category_id=category)
    if material_type:
        queryset = queryset.filter(material_type=material_type)
    if subject:
        queryset = queryset.filter(author__subject=subject)
    if grade_level:
        queryset = queryset.filter(grade_level=grade_level)
    
    # Natijalarni tartiblash
    sort_by = request.query_params.get('sort', '-created_at')
    queryset = queryset.order_by(sort_by)
    
    serializer = MaterialSerializer(queryset, many=True)
    return Response(serializer.data)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def my_materials(request):
    """Foydalanuvchining materiallari"""
    materials = Material.objects.filter(author=request.user).order_by('-created_at')
    serializer = MaterialSerializer(materials, many=True)
    return Response(serializer.data)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def material_stats(request):
    """Material statistikasi"""
    user = request.user
    
    total_materials = Material.objects.filter(author=user).count()
    total_downloads = sum(m.download_count for m in Material.objects.filter(author=user))
    avg_rating = Material.objects.filter(author=user).aggregate(
        avg_rating=models.Avg('rating')
    )['avg_rating'] or 0
    
    # Eng ko'p yuklab olingan materiallar
    popular_materials = Material.objects.filter(author=user).order_by('-download_count')[:5]
    
    return Response({
        'total_materials': total_materials,
        'total_downloads': total_downloads,
        'avg_rating': round(avg_rating, 2),
        'popular_materials': MaterialSerializer(popular_materials, many=True).data
    })


# ============ ASSIGNMENT VIEWS ============

class AssignmentListView(generics.ListCreateAPIView):
    """Topshiriqlar ro'yxati va yaratish"""
    serializer_class = AssignmentSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = Assignment.objects.filter(is_active=True)
        
        # Filtrlash
        category = self.request.query_params.get('category')
        assignment_type = self.request.query_params.get('type')
        subject = self.request.query_params.get('subject')
        grade_level = self.request.query_params.get('grade')
        teacher = self.request.query_params.get('teacher')
        
        if category:
            queryset = queryset.filter(category_id=category)
        if assignment_type:
            queryset = queryset.filter(assignment_type=assignment_type)
        if subject:
            queryset = queryset.filter(subject=subject)
        if grade_level:
            queryset = queryset.filter(grade_level=grade_level)
        if teacher:
            queryset = queryset.filter(teacher_id=teacher)
        
        return queryset.order_by('-created_at')
    
    def perform_create(self, serializer):
        serializer.save(teacher=self.request.user)


class AssignmentDetailView(generics.RetrieveUpdateDestroyAPIView):
    """Topshiriq tafsilotlari"""
    serializer_class = AssignmentSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        return Assignment.objects.filter(
            Q(is_active=True) | Q(teacher=self.request.user)
        )


class StudentSubmissionListView(generics.ListCreateAPIView):
    """O'quvchi topshirig'lari ro'yxati"""
    serializer_class = StudentSubmissionSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        assignment_id = self.request.query_params.get('assignment')
        if assignment_id:
            return StudentSubmission.objects.filter(assignment_id=assignment_id)
        return StudentSubmission.objects.none()
    
    def perform_create(self, serializer):
        serializer.save()


class StudentSubmissionDetailView(generics.RetrieveUpdateDestroyAPIView):
    """O'quvchi topshirig'i tafsilotlari"""
    serializer_class = StudentSubmissionSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        return StudentSubmission.objects.all()


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def grade_submission(request, pk):
    """O'quvchi ishini baholash"""
    submission = get_object_or_404(StudentSubmission, pk=pk)
    
    grade = request.data.get('grade')
    feedback = request.data.get('feedback', '')
    
    if not grade or not (0 <= int(grade) <= submission.assignment.max_points):
        return Response({
            'error': f'Ball 0 dan {submission.assignment.max_points} gacha bo\'lishi kerak'
        }, status=status.HTTP_400_BAD_REQUEST)
    
    submission.grade = grade
    submission.feedback = feedback
    submission.status = 'graded'
    submission.graded_by = request.user
    submission.graded_at = timezone.now()
    submission.save()
    
    return Response({
        'message': 'Ish muvaffaqiyatli baholandi',
        'submission': StudentSubmissionSerializer(submission).data
    })


# ============ VIDEO LESSON VIEWS ============

class VideoLessonListView(generics.ListCreateAPIView):
    """Video darsliklar ro'yxati"""
    serializer_class = VideoLessonSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = VideoLesson.objects.filter(is_public=True)
        
        # Filtrlash
        category = self.request.query_params.get('category')
        subject = self.request.query_params.get('subject')
        grade_level = self.request.query_params.get('grade')
        search = self.request.query_params.get('search')
        
        if category:
            queryset = queryset.filter(category_id=category)
        if subject:
            queryset = queryset.filter(subject=subject)
        if grade_level:
            queryset = queryset.filter(grade_level=grade_level)
        if search:
            queryset = queryset.filter(
                Q(title__icontains=search) |
                Q(description__icontains=search) |
                Q(tags__icontains=search)
            )
        
        return queryset.order_by('-created_at')
    
    def perform_create(self, serializer):
        serializer.save(author=self.request.user)


class VideoLessonDetailView(generics.RetrieveUpdateDestroyAPIView):
    """Video darslik tafsilotlari"""
    serializer_class = VideoLessonSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        return VideoLesson.objects.filter(
            Q(is_public=True) | Q(author=self.request.user)
        )


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def watch_video(request, pk):
    """Video ko'rish statistikasini yangilash"""
    video = get_object_or_404(VideoLesson, pk=pk)
    
    if not video.is_public and video.author != request.user:
        return Response({
            'error': 'Bu videoni ko\'rish huquqingiz yo\'q'
        }, status=status.HTTP_403_FORBIDDEN)
    
    video.view_count += 1
    video.save()
    
    return Response({
        'message': 'Video ko\'rish statistikasi yangilandi',
        'view_count': video.view_count
    })


# ============ 3D MODEL VIEWS ============

class Model3DListView(generics.ListCreateAPIView):
    """3D modellar ro'yxati"""
    serializer_class = Model3DSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        queryset = Model3D.objects.filter(is_public=True)
        
        # Filtrlash
        category = self.request.query_params.get('category')
        model_type = self.request.query_params.get('type')
        subject = self.request.query_params.get('subject')
        grade_level = self.request.query_params.get('grade')
        is_interactive = self.request.query_params.get('interactive')
        search = self.request.query_params.get('search')
        
        if category:
            queryset = queryset.filter(category_id=category)
        if model_type:
            queryset = queryset.filter(model_type=model_type)
        if subject:
            queryset = queryset.filter(subject=subject)
        if grade_level:
            queryset = queryset.filter(grade_level=grade_level)
        if is_interactive:
            queryset = queryset.filter(is_interactive=True)
        if search:
            queryset = queryset.filter(
                Q(title__icontains=search) |
                Q(description__icontains=search)
            )
        
        return queryset.order_by('-created_at')
    
    def perform_create(self, serializer):
        serializer.save(author=self.request.user)


class Model3DDetailView(generics.RetrieveUpdateDestroyAPIView):
    """3D model tafsilotlari"""
    serializer_class = Model3DSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        return Model3D.objects.filter(
            Q(is_public=True) | Q(author=self.request.user)
        )


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def download_3d_model(request, pk):
    """3D modelni yuklab olish"""
    model = get_object_or_404(Model3D, pk=pk)
    
    if not model.is_public and model.author != request.user:
        return Response({
            'error': 'Bu modelni yuklab olish huquqingiz yo\'q'
        }, status=status.HTTP_403_FORBIDDEN)
    
    try:
        # Yuklab olish statistikasini yangilash
        model.download_count += 1
        model.save()
        
        # Faylni yuklab olish
        if os.path.exists(model.model_file.path):
            with open(model.model_file.path, 'rb') as f:
                response = HttpResponse(f.read(), content_type='application/octet-stream')
                response['Content-Disposition'] = f'attachment; filename="{model.title}.{model.model_file.name.split(".")[-1]}"'
                return response
        else:
            return Response({
                'error': 'Fayl topilmadi'
            }, status=status.HTTP_404_NOT_FOUND)
            
    except Exception as e:
        return Response({
            'error': 'Faylni yuklab olishda xatolik'
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# ============ AI-ASSISTED MATERIAL CREATOR ============

_openai_client = None
_openrouter_client = None
_gemini_client = None


def _get_openai_client():
    global _openai_client
    api_key = getattr(settings, 'OPENAI_API_KEY', '')
    if not api_key:
        raise ValueError("OpenAI API kaliti topilmadi. Iltimos, .env faylida OPENAI_API_KEY ni belgilang.")
    if _openai_client is None:
        _openai_client = OpenAI(api_key=api_key)
    return _openai_client


def _get_openrouter_client():
    global _openrouter_client
    api_key = getattr(settings, 'OPENROUTER_API_KEY', '')
    if not api_key:
        raise ValueError("OpenRouter API kaliti topilmadi. Iltimos, .env faylida OPENROUTER_API_KEY ni belgilang.")
    if _openrouter_client is None:
        base_url = getattr(settings, 'OPENROUTER_BASE_URL', 'https://openrouter.ai/api/v1')
        default_headers = {}
        site_url = getattr(settings, 'OPENROUTER_SITE_URL', '')
        site_name = getattr(settings, 'OPENROUTER_SITE_NAME', '')
        if site_url:
            default_headers['HTTP-Referer'] = site_url
        if site_name:
            default_headers['X-Title'] = site_name
        _openrouter_client = OpenAI(
            api_key=api_key,
            base_url=base_url,
            default_headers=default_headers or None,
        )
    return _openrouter_client


def _generate_openai_chat(system_prompt: str, user_prompt: str, temperature: float = 0.7) -> str:
    model_name = getattr(settings, 'OPENAI_CHAT_MODEL', 'gpt-4o-mini')
    try:
        client = _get_openai_client()
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=temperature,
        )
        return response.choices[0].message.content.strip()
    except OpenAIError as exc:
        logger.exception("OpenAI chat completion failed")
        raise ValueError(f"OpenAI xatosi: {exc}") from exc
    except Exception as exc:  # noqa: BLE001
        logger.exception("OpenAI chat completion failed")
        raise ValueError(f"OpenAI xatosi: {exc}") from exc


def _generate_openrouter_chat(system_prompt: str, user_prompt: str, temperature: float = 0.7) -> str:
    configured_model = getattr(settings, 'OPENROUTER_MODEL', 'openrouter/meta/llama-3.1-8b-instruct')
    normalized = configured_model.lower().strip()
    deprecated_models = {
        'openrouter/openai/gpt-4o-mini',
        'openrouter/openai/gpt-4o-mini-2024-05',
        'openai/gpt-4o-mini',
        'gpt-4o-mini',
    }
    if normalized in deprecated_models or 'gpt-4o-mini' in normalized:
        logger.warning(
            "Deprecated OpenRouter modeli \"%s\" aniqlangan. Fallback sifatida `openai/gpt-4o` ishlatiladi. "
            "Iltimos, .env faylida OPENROUTER_MODEL qiymatini yangilang.",
            configured_model,
        )
        model_name = 'openai/gpt-4o'
    else:
        model_name = configured_model
    try:
        client = _get_openrouter_client()
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=temperature,
        )
        return response.choices[0].message.content.strip()
    except OpenAIError as exc:
        logger.exception("OpenRouter chat completion failed")
        raise ValueError(f"OpenRouter xatosi: {exc}") from exc
    except Exception as exc:  # noqa: BLE001
        logger.exception("OpenRouter chat completion failed")
        raise ValueError(f"OpenRouter xatosi: {exc}") from exc


def _generate_openai_image(prompt: str, size: str = '1024x1024') -> str:
    model_name = getattr(settings, 'OPENAI_IMAGE_MODEL', 'gpt-image-1')
    try:
        client = _get_openai_client()
        response = client.images.generate(
            model=model_name,
            prompt=prompt,
            size=size,
            quality='standard',
        )
        image_data = response.data[0].b64_json
        if not image_data:
            raise ValueError("OpenAI tasvirni qaytarmadi.")
        return image_data
    except OpenAIError as exc:
        logger.exception("OpenAI image generation failed")
        raise ValueError(f"OpenAI rasm generatsiyasi xatosi: {exc}") from exc
    except Exception as exc:  # noqa: BLE001
        logger.exception("OpenAI image generation failed")
        raise ValueError(f"OpenAI rasm generatsiyasi xatosi: {exc}") from exc


def _generate_gemini_text(prompt: str, temperature: float = 0.7) -> str:
    if genai is None:
        raise ValueError("google-genai kutubxonasi o'rnatilmagan. `pip install google-genai` buyruqini bajaring.")
    try:
        global _gemini_client
        if _gemini_client is None:
            api_key = getattr(settings, 'GOOGLE_GEMINI_API_KEY', '')
            if not api_key:
                raise ValueError("Google Gemini API kaliti topilmadi. Iltimos, .env faylida GOOGLE_GEMINI_API_KEY ni belgilang.")
            _gemini_client = genai.Client(api_key=api_key)

        configured_model = getattr(settings, 'GOOGLE_GEMINI_MODEL', 'models/gemini-1.5-flash-latest')
        model_name = configured_model if configured_model.startswith("models/") else f"models/{configured_model}"

        response = _gemini_client.models.generate_content(
            model=model_name,
            contents=[{"role": "user", "parts": [{"text": prompt}]}],
            config={"temperature": float(temperature)},
        )

        text = getattr(response, 'text', '') or ''
        if text:
            return text.strip()

        aggregates = []
        for candidate in getattr(response, 'candidates', []) or []:
            for part in getattr(getattr(candidate, 'content', None), 'parts', []) or []:
                part_text = getattr(part, 'text', '')
                if part_text:
                    aggregates.append(part_text)

        if aggregates:
            return "\n\n".join(aggregates).strip()

        finish_reasons = [
            getattr(candidate, 'finish_reason', '')
            for candidate in getattr(response, 'candidates', []) or []
        ]
        if finish_reasons:
            raise ValueError(f"Gemini javobi to'liq emas: {', '.join(filter(None, finish_reasons))}.")

        raise ValueError("Gemini bo'sh javob qaytardi.")
    except Exception as exc:  # noqa: BLE001
        logger.exception("Gemini generate_content failed")
        raise ValueError(f"Gemini xatosi: {exc}") from exc


def _generate_ai_text(provider: str, system_prompt: str, user_prompt: str, temperature: float = 0.7) -> str:
    provider = (provider or 'openai').lower()
    if provider == 'gemini':
        combined_prompt = f"{system_prompt}\n\nFoydalanuvchi so'rovi:\n{user_prompt}\n\nNatijani tartibli va tushunarli shaklda taqdim et."
        return _generate_gemini_text(combined_prompt, temperature=temperature)
    if provider == 'openrouter':
        return _generate_openrouter_chat(system_prompt, user_prompt, temperature=temperature)
    return _generate_openai_chat(system_prompt, user_prompt, temperature=temperature)


def _parse_presentation_slides(text: str) -> list[dict]:
    if not text:
        return []

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    slides: list[dict] = []
    current_title = None
    current_bullets: list[str] = []

    slide_break_pattern = re.compile(r'^(?:slide\s*)?\d+[\).:-]\s*', re.IGNORECASE)

    for line in lines:
        normalized = line.replace('—', '-').strip()
        if slide_break_pattern.match(normalized):
            if current_title or current_bullets:
                slides.append({
                    'title': current_title or "Slayd",
                    'bullets': current_bullets[:] or ["Mazmun kiritilmagan"],
                })
            current_title = slide_break_pattern.sub('', normalized).strip() or "Slayd"
            current_bullets = []
            continue

        if current_title is None:
            current_title = normalized
            continue

        bullet = re.sub(r'^[-•*]\s*', '', normalized).strip()
        current_bullets.append(bullet or normalized)

    if current_title or current_bullets:
        slides.append({
            'title': current_title or "Slayd",
            'bullets': current_bullets[:] or ["Mazmun kiritilmagan"],
        })

    if not slides:
        slides.append({
            'title': "Slayd mazmuni",
            'bullets': [text.strip()],
        })

    return slides


def _build_presentation_stream(topic: str, content: str) -> tuple[io.BytesIO, str]:
    slides = _parse_presentation_slides(content)

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = topic or "AI taqdimot"
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Ustoziya AI tomonidan yaratilgan taqdimot"

    for slide_data in slides:
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data['title']
        text_frame = slide.shapes.placeholders[1].text_frame
        bullets = [b for b in slide_data['bullets'] if b]
        if not bullets:
            bullets = ["Mazmun mavjud emas"]

        text_frame.text = bullets[0]
        text_frame.paragraphs[0].level = 0
        for bullet in bullets[1:]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0

    output = io.BytesIO()
    presentation.save(output)
    output.seek(0)

    filename_root = slugify(topic) or "taqdimot"
    return output, f"{filename_root}.pptx"


def _generate_ai_image(provider: str, prompt: str, size: str) -> dict:
    provider = (provider or 'openai').lower()
    if provider == 'gemini':
        description_prompt = (
            "Siz tajribali grafik dizayn bo'yicha sun'iy intellektsiz. "
            "Foydalanuvchi ta'riflagan tasvir uchun batafsil dizayn tavsifi yarating. "
            "Natijada quyidagi tuzilmani saqlang:\n"
            "1. Qisqa konsepsiya tavsifi\n"
            "2. Kompozitsiya va asosiy obyektlar\n"
            "3. Rang palitrasi va yorug'lik\n"
            "4. Qo'shimcha detal va teksturalar\n"
            "5. Agar kerak bo'lsa, matn elementlari\n\n"
            f"Foydalanuvchi so'rovi: {prompt}"
        )
        description = _generate_gemini_text(description_prompt, temperature=0.6)
        return {'description': description}

    image_base64 = _generate_openai_image(prompt, size=size)
    return {'image_base64': image_base64}


def _build_ai_cards(request):
    return [
        {
            'title': "Taqdimot tayyorlash",
            'slug': 'presentation',
            'emoji': '🧾',
            'description': "Slaydlar strukturasini tez tayyorlash va matnlarni AI yordamida olish.",
            'url': reverse('materials:presentation'),
            'background': 'linear-gradient(135deg, #ff9f43, #ff6f20)',
        },
        {
            'title': "Interaktiv metodlar",
            'slug': 'interactive',
            'emoji': '🧠',
            'description': "Darslar uchun yangi metod va faol mashg'ulotlar g'oyalari.",
            'url': reverse('materials:interactive'),
            'background': 'linear-gradient(135deg, #1abc9c, #109177)',
        },
        {
            'title': "Rasm tayyorlash",
            'slug': 'image',
            'emoji': '🖼️',
            'description': "Ko'rgazmali tasvirlar yoki infografikalarni yaratish.",
            'url': reverse('materials:image'),
            'background': 'linear-gradient(135deg, #3498db, #2c82c9)',
        },
        {
            'title': "Qo'shimcha o'quv manbalari",
            'slug': 'resources',
            'emoji': '📚',
            'description': "Kitoblar, maqolalar va onlayn resurslar bo'yicha AI tavsiyalari.",
            'url': reverse('materials:resources'),
            'background': 'linear-gradient(135deg, #f6d365, #fda085)',
        },
    ]


@login_required
def material_ai_categories(request):
    context = {
        'cards': _build_ai_cards(request),
    }
    return render(request, 'materials/categories.html', context)


def _process_text_generation(request, *, template_name: str, system_prompt: str, hero: dict, temperature: float = 0.7):
    try:
        default_provider = _resolve_text_provider()
    except ValueError as exc:
        context = {
            'hero': hero,
            'prompt_value': '',
            'provider_label': "API kaliti topilmadi",
            'error': str(exc),
        }
        return render(request, template_name, context)

    context = {
        'hero': hero,
        'prompt_value': '',
        'provider_label': _get_provider_label(default_provider, TEXT_PROVIDER_LABELS),
    }

    if request.method == 'POST':
        prompt = request.POST.get('prompt', '').strip()
        provider = default_provider

        export_ppt = request.POST.get('export_ppt') == '1'
        existing_result = request.POST.get('existing_result', '').strip()

        context['prompt_value'] = prompt

        if not prompt:
            context['error'] = "Iltimos, so'rov matnini kiriting."
        else:
            if export_ppt:
                if not existing_result:
                    try:
                        existing_result = _generate_ai_text(provider, system_prompt, prompt, temperature=temperature)
                    except ValueError as exc:
                        context['error'] = str(exc)
                if existing_result and 'error' not in context:
                    try:
                        stream, filename = _build_presentation_stream(prompt, existing_result)
                        return FileResponse(stream, as_attachment=True, filename=filename)
                    except ValueError as exc:
                        context['error'] = str(exc)
                    except Exception as exc:  # noqa: BLE001
                        logger.exception("PPT generation failed")
                        context['error'] = f"PPT yaratishda kutilmagan xatolik: {exc}"
                context['result'] = existing_result
            else:
                try:
                    result = _generate_ai_text(provider, system_prompt, prompt, temperature=temperature)
                    context['result'] = result
                except ValueError as exc:
                    context['error'] = str(exc)

        if 'result' in context:
            context['provider_label'] = _get_provider_label(provider, TEXT_PROVIDER_LABELS)

    return render(request, template_name, context)


@ensure_csrf_cookie
@login_required
def material_presentation_view(request):
    system_prompt = (
        "Siz ta'lim sohasida ishlaydigan professional taqdimot dizaynerisiz. "
        "Berilgan mavzu bo'yicha 5-7 slayddan iborat taqdimot rejasini yarating. "
        "Har bir slayd uchun sarlavha va 3-5 ta asosiy bullet yozing. "
        "Zarur bo'lsa, ko'rgazmali elementlar va interaktiv savollar uchun tavsiyalar qo'shing."
    )
    hero = {
        'emoji': '🧾',
        'title': "AI yordamida taqdimot",
        'subtitle': "Mavzuni kiriting va AI slayd reja hamda mazmunini taklif etsin.",
    }
    return _process_text_generation(
        request,
        template_name='materials/presentation.html',
        system_prompt=system_prompt,
        hero=hero,
    )


@ensure_csrf_cookie
@login_required
def material_interactive_view(request):
    system_prompt = (
        "Siz ijodkor metodist va sinf rahbarisiz. "
        "Berilgan mavzu yoki yosh toifasi uchun interaktiv dars metodi tavsiyasini ishlab chiqing. "
        "Metodning maqsadi, zarur materiallar, bosqichma-bosqich o'tkazish tartibi va yakuniy refleksiya savollarini ko'rsating."
    )
    hero = {
        'emoji': '🧠',
        'title': "AI bilan interaktiv metod",
        'subtitle': "AI siz uchun qadam-baqadam faol dars ssenariysini taklif etadi.",
    }
    return _process_text_generation(
        request,
        template_name='materials/interactive.html',
        system_prompt=system_prompt,
        hero=hero,
    )


@ensure_csrf_cookie
@login_required
def material_resources_view(request):
    system_prompt = (
        "Siz tajribali kutubxonachi va ta'lim metodologi sifatida ishlaysiz. "
        "Berilgan fan, mavzu yoki ko'nikma uchun 5-7 ta qo'shimcha ta'lim manbalarini taklif eting. "
        "Har bir manba uchun qisqa tavsif, foydalanish maqsadi va havola (agar mavjud bo'lsa) ko'rsatilishi kerak."
    )
    hero = {
        'emoji': '📚',
        'title': "AI tavsiya qilgan manbalar",
        'subtitle': "O'quvchilar va o'qituvchilar uchun eng mos qo'shimcha resurslarni toping.",
    }
    return _process_text_generation(
        request,
        template_name='materials/resources.html',
        system_prompt=system_prompt,
        hero=hero,
        temperature=0.6,
    )


@ensure_csrf_cookie
@login_required
def material_image_view(request):
    hero = {
        'emoji': '🖼️',
        'title': "AI rasm va vizual kontent",
        'subtitle': "Tasvirni tasvirlab bering – AI uni generatsiya qiladi yoki dizayn tavsiyasi beradi.",
    }
    allowed_sizes = ['512x512', '768x768', '1024x1024']
    try:
        provider = _resolve_image_provider()
    except ValueError as exc:
        context = {
            'hero': hero,
            'prompt_value': '',
            'size': '1024x1024',
            'allowed_sizes': allowed_sizes,
            'provider_label': "API kaliti topilmadi",
            'error': str(exc),
        }
        return render(request, 'materials/image.html', context)

    context = {
        'hero': hero,
        'prompt_value': '',
        'size': '1024x1024',
        'allowed_sizes': allowed_sizes,
        'provider_label': _get_provider_label(provider, IMAGE_PROVIDER_LABELS),
        'selected_provider': provider,
    }

    if request.method == 'POST':
        prompt = request.POST.get('prompt', '').strip()
        size = request.POST.get('size', '1024x1024')

        if size not in allowed_sizes:
            size = '1024x1024'

        context['prompt_value'] = prompt
        context['size'] = size

        if not prompt:
            context['error'] = "Iltimos, rasm tavsifini kiriting."
        else:
            try:
                context.pop('image_data_url', None)
                context.pop('image_base64', None)
                context.pop('description', None)
                result = _generate_ai_image(context['selected_provider'], prompt, size)
                context.update(result)
                context['provider_label'] = _get_provider_label(context['selected_provider'], IMAGE_PROVIDER_LABELS)
                if 'image_base64' in result:
                    context['image_data_url'] = f"data:image/png;base64,{result['image_base64']}"
            except ValueError as exc:
                context['error'] = str(exc)

    return render(request, 'materials/image.html', context)